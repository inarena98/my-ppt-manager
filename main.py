import os
import pandas as pd
from pptx import Presentation

# ==========================================
# [기능 1] PPT 목차 추출 -> 엑셀 파일 생성
# ==========================================
def extract_toc_to_excel(ppt_path, excel_path):
    """
    원본 PPT 파일에서 슬라이드 번호와 제목을 추출하여 엑셀 가이드라인을 생성합니다.
    """
    if not os.path.exists(ppt_path):
        print(f"❌ 오류: 원본 PPT 파일 '{ppt_path}'을 찾을 수 없습니다.")
        return

    prs = Presentation(ppt_path)
    data = []

    print(f"🔍 '{ppt_path}'에서 목차 추출 중...")
    for i, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        
        # 슬라이드 제목 구역 탐색
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_title = slide.shapes.title.text.strip()
        
        # 제목 구역이 없으면 첫 번째 발견되는 텍스트 상자의 첫 줄 사용
        if not slide_title:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    slide_title = shape.text.split('\n')[0].strip()
                    break
        
        if not slide_title:
            slide_title = f"제목 없는 슬라이드 {i}"

        # 엑셀에 작성할 초기 데이터 구조 생성
        data.append({
            "순서": i,
            "슬라이드 주제(목차)": slide_title,
            "개별_PPT_파일명": f"slide_{i:02d}.pptx", # 매핑할 기본 파일명 추천
            "작성률(%)": 0 # 사용자가 직접 100으로 입력할 공간
        })

    # 엑셀 저장
    df = pd.DataFrame(data)
    df.to_excel(excel_path, index=False)
    print(f"✅ 목차 엑셀 생성 완료: '{excel_path}'")
    print("👉 생성된 엑셀 파일을 열어 '개별_PPT_파일명'과 '작성률(%)'을 업데이트하세요.")


# ==========================================
# [기능 2] 엑셀 기준 -> 작성률 100% PPT 취합
# ==========================================
def merge_ppts_by_excel(excel_path, ppt_folder, output_path):
    """
    엑셀 파일을 읽어 작성률이 100%인 개별 PPT 파일들을 순서대로 하나의 파일로 취합합니다.
    """
    if not os.path.exists(excel_path):
        print(f"❌ 오류: 엑셀 파일 '{excel_path}'을 찾을 수 없습니다.")
        return

    # 엑셀 데이터 로드
    df = pd.read_excel(excel_path)
    
    # 필수 열 확인
    required_cols = ["순서", "개별_PPT_파일명", "작성률(%)"]
    if not all(col in df.columns for col in required_cols):
        print(f"❌ 오류: 엑셀 파일에 {required_cols} 열이 모두 존재해야 합니다.")
        return

    # 조건 필터링: 작성률이 100인 행만 필터링 후 순서대로 정렬
    target_files = df[df["작성률(%)"] == 100].sort_values(by="순서")

    if target_files.empty:
        print("⚠️ 취합할 대상(작성률 100%) PPT 파일이 엑셀에 존재하지 않습니다.")
        return

    # 새로운 빈 프레젠테이션 생성 (취합본이 될 파일)
    merged_prs = Presentation()
    # 기본 생성 시 생기는 첫 번째 빈 슬라이드 제거 (선택 사항)
    if merged_prs.slides:
        rId = merged_prs.slides._sldIdLst[0].rId
        merged_prs.part.drop_rel(rId)
        del merged_prs.slides._sldIdLst[0]

    print(f"\n📂 PPT 취합 시작 (총 {len(target_files)}개 대상)...")
    
    merged_count = 0
    for idx, row in target_files.iterrows():
        filename = str(row["개별_PPT_파일명"]).strip()
        slide_title = row["슬라이드 주제(목차)"]
        file_path = os.path.join(ppt_folder, filename)

        if not os.path.exists(file_path):
            print(f"⚠️ 경고: '{filename}' 파일이 폴더 내에 없습니다. (목차: {slide_title}) - 건너뜁니다.")
            continue

        print(f"➡️ [{row['순서']}번] '{filename}' 병합 중... ({slide_title})")
        
        # 개별 PPT 불러오기
        sub_prs = Presentation(file_path)
        
        # 슬라이드 복사 및 추가
        for slide in sub_prs.slides:
            # 원본 슬라이드의 레이아웃을 빈 프레젠테이션에 추가
            slide_layout = merged_prs.slide_layouts[5] # 빈 슬라이드 레이아웃 혹은 원본에 대응하는 레이아웃 선택
            new_slide = merged_prs.slides.add_slide(slide_layout)
            
            # 도형, 텍스트, 이미지 요소들을 복제
            for shape in slide.shapes:
                # 간단한 텍스트 및 기본 형태의 엘리먼트 복제 (python-pptx 한계 상 단순 복제 예시)
                # 서식을 완벽하게 보존하기 위해선 윈도우 환경의 win32com 제어가 유리합니다.
                if shape.has_text_frame:
                    txBox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                    tf = txBox.text_frame
                    tf.text = shape.text_frame.text
                    
        merged_count += 1

    if merged_count > 0:
        merged_prs.save(output_path)
        print(f"\n🎉 취합 완료! 최종 파일 저장됨: '{output_path}'")
    else:
        print("\n❌ 취합된 파일이 없습니다.")


# ==========================================
# 프로그램 실행 진입점
# ==========================================
if __name__ == "__main__":
    # 사용할 파일명 및 폴더 정의
    ORIGINAL_PPT = "master_template.pptx"      # 목차를 뽑아낼 최초 원본 파일
    GUIDE_EXCEL = "ppt_merge_guide.xlsx"       # 추출하여 생성할 가이드 엑셀 파일
    PPT_SOURCE_DIR = "./individual_ppts"       # 개별 PPT 조각들이 모여있는 폴더
    FINAL_MERGED_PPT = "final_output.pptx"     # 최종 취합될 파일명

    # 실행할 모드 선택 (콘솔 인터페이스)
    print("====== PPT 목차 관리 시스템 ======")
    print("1: 원본 PPT에서 목차 추출하여 엑셀 만들기")
    print("2: 엑셀 기준 작성률 100% PPT 취합하기")
    choice = input("원하는 작업 번호를 입력하세요 (1 또는 2): ").strip()

    if choice == "1":
        extract_toc_to_excel(ORIGINAL_PPT, GUIDE_EXCEL)
    elif choice == "2":
        # 폴더가 없으면 생성
        if not os.path.exists(PPT_SOURCE_DIR):
            os.makedirs(PPT_SOURCE_DIR)
            print(f"📂 '{PPT_SOURCE_DIR}' 폴더가 없어 생성했습니다. 개별 PPT들을 이 폴더에 넣어주세요.")
        merge_ppts_by_excel(GUIDE_EXCEL, PPT_SOURCE_DIR, FINAL_MERGED_PPT)
    else:
        print("❌ 잘못된 입력입니다. 프로그램을 종료합니다.")
